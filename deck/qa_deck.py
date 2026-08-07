"""QA for the generated deck: package integrity, geometry, and text fit.

    .venv/bin/python deck/qa_deck.py ML_Bubble_2026_Readmission_Risk.pptx

Stands in for the skill's validate.py, which needs Python 3.10+ (this machine has
3.9) and for image-level visual QA, which needs LibreOffice (not installed).
Checks the defects that actually reach a viewer: parts that don't resolve, shapes
off the slide or crowding the margin, overlapping text, and text too long for its
box.
"""

from __future__ import annotations

import sys
import zipfile
from collections import defaultdict
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.util import Emu

MARGIN_MIN = 0.5      # inches from slide edge
GAP_MIN = 0.0         # overlap tolerance in inches (text boxes only)
# Rough average glyph width as a fraction of point size. Calibri/Cambria sit
# near 0.48; deliberately optimistic so the checker flags only real problems.
CHAR_W = 0.47
LINE_H = 1.22         # line height as a multiple of font size


def emu_in(v) -> float:
    return Emu(v).inches if v is not None else 0.0


def check_package(path: Path) -> list[str]:
    """Every relationship target and content-type override must resolve."""
    problems = []
    with zipfile.ZipFile(path) as z:
        names = set(z.namelist())
        bad = z.testzip()
        if bad:
            problems.append(f"corrupt zip entry: {bad}")

        ct = ET.fromstring(z.read("[Content_Types].xml"))
        ns = "{http://schemas.openxmlformats.org/package/2006/content-types}"
        for ov in ct.findall(f"{ns}Override"):
            part = ov.get("PartName", "").lstrip("/")
            if part not in names:
                problems.append(f"content-type override for missing part: {part}")

        rns = "{http://schemas.openxmlformats.org/package/2006/relationships}"
        for rels in [n for n in names if n.endswith(".rels")]:
            base = str(Path(rels).parent.parent)
            for rel in ET.fromstring(z.read(rels)).findall(f"{rns}Relationship"):
                if rel.get("TargetMode") == "External":
                    continue
                target = rel.get("Target", "")
                resolved = str((Path(base) / target).resolve()).lstrip("/")
                # normalise against the archive root
                resolved = resolved.split(str(Path.cwd()).lstrip("/") + "/")[-1]
                candidates = {resolved, str(Path(base) / target).replace("../", "")}
                if not (candidates & names) and not any(
                    n.endswith(Path(target).name) for n in names
                ):
                    problems.append(f"{rels}: unresolved target {target}")

        charts = [n for n in names if n.startswith("ppt/charts/") and n.endswith(".xml")]
        for c in charts:
            body = z.read(c).decode("utf8", "ignore")
            if "secondaryValAxis" in body and body.count("<c:valAx>") < 2:
                problems.append(f"{c}: secondary axis declared without two valAx entries")
        print(f"  package: {len(names)} parts, {len(charts)} chart(s)")
    return problems


def text_of(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs)


def max_font(shape) -> float:
    sizes = [
        r.font.size.pt
        for p in shape.text_frame.paragraphs
        for r in p.runs
        if r.font.size is not None
    ]
    return max(sizes) if sizes else 18.0


def check_slides(path: Path) -> tuple[list[str], list[str]]:
    prs = Presentation(str(path))
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    problems, warnings = [], []

    for idx, slide in enumerate(prs.slides, 1):
        boxes = []
        for shape in slide.shapes:
            x, y = emu_in(shape.left), emu_in(shape.top)
            w, h = emu_in(shape.width), emu_in(shape.height)
            txt = text_of(shape).strip()

            # --- bounds -------------------------------------------------- #
            if x < -0.01 or y < -0.01 or x + w > sw + 0.01 or y + h > sh + 0.01:
                # decorative circles are intentionally bled off the edge
                if txt:
                    problems.append(
                        f"slide {idx}: '{txt[:38]}' outside slide "
                        f"({x:.2f},{y:.2f} {w:.2f}x{h:.2f})"
                    )
            if txt and (x < MARGIN_MIN - 0.01 or x + w > sw - MARGIN_MIN + 0.01):
                warnings.append(
                    f"slide {idx}: '{txt[:32]}' within {MARGIN_MIN}\" of a side edge"
                )
            if txt and (y < 0.3 or y + h > sh - 0.28):
                warnings.append(f"slide {idx}: '{txt[:32]}' close to top/bottom edge")

            # --- text fit ------------------------------------------------ #
            if txt:
                fs = max_font(shape)
                usable = max(w - 0.1, 0.4)
                lines = 0
                for para in txt.split("\n"):
                    chars_per_line = max(int(usable * 72 / (fs * CHAR_W)), 1)
                    lines += max(1, -(-len(para) // chars_per_line))
                needed = lines * fs * LINE_H / 72
                if needed > h * 1.12:
                    problems.append(
                        f"slide {idx}: text overflow — '{txt[:38]}' needs "
                        f"~{needed:.2f}\" in a {h:.2f}\" box ({fs:.0f}pt, {lines} lines)"
                    )
                boxes.append((x, y, w, h, txt, fs))

        # --- overlap between text-bearing shapes ------------------------- #
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                ax, ay, aw, ah, at, _ = boxes[i]
                bx, by, bw, bh, bt, _ = boxes[j]
                ox = min(ax + aw, bx + bw) - max(ax, bx)
                oy = min(ay + ah, by + bh) - max(ay, by)
                if ox > GAP_MIN + 0.06 and oy > GAP_MIN + 0.06:
                    area = ox * oy
                    if area > 0.06:
                        warnings.append(
                            f"slide {idx}: text overlap {area:.2f} sq\" — "
                            f"'{at[:24]}' / '{bt[:24]}'"
                        )
    return problems, warnings


def dump_text(path: Path) -> None:
    prs = Presentation(str(path))
    placeholders = ("lorem", "ipsum", "todo", "[insert", "xxx")
    hits = []
    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            t = text_of(shape).lower()
            for p in placeholders:
                if p in t:
                    hits.append(f"slide {idx}: placeholder '{p}'")
    print(f"  content: {len(prs.slides)} slides, {len(hits)} placeholder hit(s)")
    for h in hits:
        print(f"    {h}")


def main() -> int:
    path = Path(sys.argv[1] if len(sys.argv) > 1 else "ML_Bubble_2026_Readmission_Risk.pptx")
    print(f"QA {path.name}")

    pkg = check_package(path)
    dump_text(path)
    problems, warnings = check_slides(path)
    problems = pkg + problems

    print(f"\n  {len(problems)} problem(s), {len(warnings)} warning(s)")
    for p in problems:
        print(f"  FAIL  {p}")
    for w in warnings:
        print(f"  warn  {w}")
    return 1 if problems else 0


if __name__ == "__main__":
    sys.exit(main())
