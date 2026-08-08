"""Render the generated deck to HTML so its layout can be inspected visually.

    .venv/bin/python deck/preview.py ML_Bubble_2026_Readmission_Risk.pptx

LibreOffice is not available on this machine, so slides cannot be rasterised the
usual way. This reads the *written* .pptx with python-pptx — not the generator
source — and lays every shape out at its real position and size, so what appears
here is what the file actually contains. Fonts are approximated by the browser,
so treat text fit as indicative; positions, sizes and overlaps are exact.
"""

from __future__ import annotations

import base64
import html
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

SCALE = 96  # px per inch


def inches(v) -> float:
    return Emu(v).inches if v is not None else 0.0


def solid_fill(shape) -> str | None:
    try:
        f = shape.fill
        if f.type is not None and f.type == 1:  # MSO_FILL.SOLID
            return f"#{f.fore_color.rgb}"
    except Exception:
        pass
    return None


def para_html(shape) -> str:
    out = []
    for p in shape.text_frame.paragraphs:
        runs = []
        for r in p.runs:
            style = []
            if r.font.size:
                style.append(f"font-size:{r.font.size.pt * SCALE / 72:.1f}px")
            if r.font.bold:
                style.append("font-weight:700")
            if r.font.italic:
                style.append("font-style:italic")
            try:
                if r.font.color and r.font.color.rgb:
                    style.append(f"color:#{r.font.color.rgb}")
            except Exception:
                pass
            if r.font.name:
                style.append(f"font-family:'{r.font.name}',serif")
            runs.append(f'<span style="{";".join(style)}">{html.escape(r.text)}</span>')
        bullet = "• " if p.level == 0 and _has_bullet(p) else ""
        out.append(f'<div class="p">{bullet}{"".join(runs) or "&nbsp;"}</div>')
    return "".join(out)


def _has_bullet(p) -> bool:
    xml = p._p.xml
    return "buChar" in xml or "buAutoNum" in xml


def slide_html(slide, idx, sw, sh, media: dict) -> str:
    bg = "#FFFFFF"
    try:
        if slide.background.fill.type == 1:
            bg = f"#{slide.background.fill.fore_color.rgb}"
    except Exception:
        pass

    parts = [
        f'<div class="slide" style="width:{sw*SCALE}px;height:{sh*SCALE}px;background:{bg}">',
        f'<div class="num">{idx}</div>',
    ]

    for shape in slide.shapes:
        x, y = inches(shape.left) * SCALE, inches(shape.top) * SCALE
        w, h = inches(shape.width) * SCALE, inches(shape.height) * SCALE
        base = f"left:{x:.1f}px;top:{y:.1f}px;width:{w:.1f}px;height:{h:.1f}px"

        if shape.shape_type == 13 and (idx, shape.shape_id) in media:  # PICTURE
            parts.append(f'<img class="sh" style="{base}" src="{media[(idx, shape.shape_id)]}">')
            continue
        if shape.has_chart:
            parts.append(
                f'<div class="sh chart" style="{base}">native chart<br>'
                f"{shape.chart.chart_type}</div>"
            )
            continue
        if shape.has_table:
            rows = []
            for r in shape.table.rows:
                cells = "".join(
                    f"<td>{html.escape(c.text)}</td>" for c in r.cells
                )
                rows.append(f"<tr>{cells}</tr>")
            parts.append(
                f'<table class="sh tbl" style="{base}">{"".join(rows)}</table>'
            )
            continue

        fill = solid_fill(shape)
        # The specific geometry lives on auto_shape_type; shape_type just says
        # AUTO_SHAPE for every rounded rectangle and ellipse alike.
        geom = ""
        try:
            geom = str(shape.auto_shape_type)
        except Exception:
            pass
        radius = "border-radius:8px;" if "ROUNDED_RECT" in geom else ""
        if "OVAL" in geom or "ELLIPSE" in geom:
            radius = "border-radius:50%;"
        style = base + ";" + radius + (f"background:{fill};" if fill else "")
        body = para_html(shape) if shape.has_text_frame else ""
        cls = "sh box" if fill else "sh"
        parts.append(f'<div class="{cls}" style="{style}">{body}</div>')

    parts.append("</div>")
    return "".join(parts)


def main() -> int:
    src = Path(sys.argv[1] if len(sys.argv) > 1 else "ML_Bubble_2026_Readmission_Risk.pptx")
    # Optional 1-indexed slide filter, e.g. `preview.py deck.pptx 11,13,15`.
    # Embedded images are large, so a whole-deck page can overwhelm a renderer;
    # this keeps a focused review light enough to load.
    wanted = None
    out_name = "preview.html"
    if len(sys.argv) > 2:
        wanted = {int(n) for n in sys.argv[2].split(",")}
        out_name = "preview_subset.html"

    prs = Presentation(str(src))
    sw, sh = prs.slide_width.inches, prs.slide_height.inches

    # shape_id is unique per slide, not per presentation — key on both.
    media = {}
    for si, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == 13:
                try:
                    img = shape.image
                    media[(si, shape.shape_id)] = (
                        f"data:{img.content_type};base64,"
                        + base64.b64encode(img.blob).decode()
                    )
                except Exception:
                    pass

    body = "".join(
        slide_html(s, i, sw, sh, media)
        for i, s in enumerate(prs.slides, 1)
        if wanted is None or i in wanted
    )
    doc = f"""<!doctype html><meta charset="utf-8">
<title>Deck preview — {html.escape(src.name)}</title>
<style>
 body {{ background:#3a4750; margin:0; padding:24px; font-family:Calibri,Arial,sans-serif; }}
 .slide {{ position:relative; margin:0 auto 28px; box-shadow:0 6px 22px rgba(0,0,0,.45);
           overflow:hidden; }}
 .sh {{ position:absolute; overflow:visible; }}
 .p {{ line-height:1.22; }}
 .chart {{ background:#dde6ea; color:#2b4a56; display:flex; align-items:center;
           justify-content:center; font-size:13px; text-align:center; }}
 .tbl {{ border-collapse:collapse; font-size:12px; }}
 .tbl td {{ border:1px solid #dce7eb; padding:2px 6px; }}
 .num {{ position:absolute; right:6px; top:4px; font-size:11px; color:#b9c6cc; z-index:9; }}
</style>
{body}
"""
    out = Path("deck") / out_name
    out.write_text(doc)
    shown = len(wanted) if wanted else len(prs.slides)
    print(f"wrote {out} — {shown} of {len(prs.slides)} slides, {len(media)} embedded image(s)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
